"""Sprint A round 2 -- docx / pptx / xlsx loaders.

Each test builds a tiny document in-memory (no fixtures on disk) so the
test suite stays self-contained + deterministic. We verify:

- The loader is registered in LOADERS for the right suffix
- Title + text come back non-empty
- Structural markers (headings, tables, slide sections) survive the
  round-trip so the downstream chunker still has something to split on
"""
from __future__ import annotations

from pathlib import Path

import pytest

from rag.ingest.loaders import LOADERS, load_docx, load_file, load_pptx, load_xlsx


# ---------- docx ----------

@pytest.fixture
def docx_path(tmp_path: Path) -> Path:
    import docx
    doc = docx.Document()
    doc.core_properties.title = "Sample Doc Title"
    doc.add_heading("Top-level heading", level=1)
    doc.add_paragraph("First paragraph under heading 1.")
    doc.add_heading("Sub heading", level=2)
    doc.add_paragraph("A body paragraph with some content.")
    tbl = doc.add_table(rows=2, cols=3)
    tbl.rows[0].cells[0].text = "col a"
    tbl.rows[0].cells[1].text = "col b"
    tbl.rows[0].cells[2].text = "col c"
    tbl.rows[1].cells[0].text = "1"
    tbl.rows[1].cells[1].text = "2"
    tbl.rows[1].cells[2].text = "3"
    doc.add_paragraph("Final paragraph after the table.")
    p = tmp_path / "sample.docx"
    doc.save(str(p))
    return p


def test_docx_loader_registered():
    assert ".docx" in LOADERS
    assert LOADERS[".docx"] is load_docx


def test_docx_loader_extracts_title_and_structure(docx_path: Path):
    title, text = load_docx(docx_path)
    assert title == "Sample Doc Title"
    # Heading 1 -> "# ..."
    assert "# Top-level heading" in text
    # Heading 2 -> "## ..."
    assert "## Sub heading" in text
    # Body paragraphs preserved
    assert "First paragraph under heading 1." in text
    assert "Final paragraph after the table." in text
    # Table rendered as markdown pipe table
    assert "| col a | col b | col c |" in text
    assert "| 1 | 2 | 3 |" in text


def test_docx_preserves_block_order(docx_path: Path):
    _, text = load_docx(docx_path)
    # Paragraph BEFORE table, final paragraph AFTER table.
    i_pre = text.find("First paragraph under heading 1.")
    i_table = text.find("| col a | col b | col c |")
    i_post = text.find("Final paragraph after the table.")
    assert 0 <= i_pre < i_table < i_post


def test_load_file_dispatches_docx(docx_path: Path):
    title, text = load_file(docx_path)
    assert title and text


# ---------- pptx ----------

@pytest.fixture
def pptx_path(tmp_path: Path) -> Path:
    from pptx import Presentation
    prs = Presentation()
    prs.core_properties.title = "Quarterly Review"

    # Slide 1: title layout
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction"
    # Title layout has a subtitle placeholder
    slide1.placeholders[1].text = "Intro subtitle body text"

    # Slide 2: blank layout with manually added text
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide2.shapes.title.text = "Key Numbers"
    # Speaker notes
    slide2.notes_slide.notes_text_frame.text = "Remember to mention the 30% growth."

    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    return p


def test_pptx_loader_registered():
    assert ".pptx" in LOADERS
    assert LOADERS[".pptx"] is load_pptx


def test_pptx_loader_emits_slide_sections(pptx_path: Path):
    title, text = load_pptx(pptx_path)
    assert title == "Quarterly Review"
    # Slide section headers
    assert "## Slide 1: Introduction" in text
    assert "## Slide 2: Key Numbers" in text
    # Body + notes preserved
    assert "Intro subtitle body text" in text
    assert "> Notes: Remember to mention the 30% growth." in text


def test_pptx_section_order(pptx_path: Path):
    _, text = load_pptx(pptx_path)
    assert text.find("## Slide 1") < text.find("## Slide 2")


# ---------- xlsx ----------

@pytest.fixture
def xlsx_path(tmp_path: Path) -> Path:
    from openpyxl import Workbook
    wb = Workbook()
    wb.properties.title = "Orders Sheet"
    ws = wb.active
    ws.title = "Q1"
    ws.append(["order_id", "sku", "qty"])
    ws.append(["O-1", "ABC", 2])
    ws.append(["O-2", "XYZ", 5])

    ws2 = wb.create_sheet("Empty")  # should be skipped
    ws3 = wb.create_sheet("Q2")
    ws3.append(["order_id", "sku", "qty"])
    ws3.append(["O-3", "LMN", 1])

    p = tmp_path / "sales.xlsx"
    wb.save(str(p))
    return p


def test_xlsx_loader_registered():
    assert ".xlsx" in LOADERS
    assert LOADERS[".xlsx"] is load_xlsx


def test_xlsx_loader_emits_sheet_sections(xlsx_path: Path):
    title, text = load_xlsx(xlsx_path)
    assert title == "Orders Sheet"
    assert "## Sheet: Q1" in text
    assert "## Sheet: Q2" in text
    # Empty sheet should NOT appear
    assert "## Sheet: Empty" not in text


def test_xlsx_renders_markdown_table(xlsx_path: Path):
    _, text = load_xlsx(xlsx_path)
    # Header row
    assert "| order_id | sku | qty |" in text
    # Body rows
    assert "| O-1 | ABC | 2 |" in text
    assert "| O-3 | LMN | 1 |" in text
    # Separator row
    assert "| --- | --- | --- |" in text


def test_load_file_dispatches_xlsx(xlsx_path: Path):
    title, text = load_file(xlsx_path)
    assert title and text and "## Sheet" in text
