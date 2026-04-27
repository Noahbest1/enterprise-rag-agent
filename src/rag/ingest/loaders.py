"""File loaders. Each returns (title, markdown_text).

Keep markdown structure (headings, lists, tables) so the chunker can split on it.
PDFs with a real text layer skip OCR; scanned PDFs fall through to the OCR loader.
"""
from __future__ import annotations

import re
from pathlib import Path


def load_text(path: Path) -> tuple[str, str]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return path.stem, text


def load_markdown(path: Path) -> tuple[str, str]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    title = path.stem
    m = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
    if m:
        title = m.group(1).strip()
    return title, text


def load_html(path: Path) -> tuple[str, str]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(raw, "html.parser")
        title = soup.title.string.strip() if soup.title and soup.title.string else path.stem
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text("\n", strip=True)
        return title, text
    except Exception:
        stripped = re.sub(r"<[^>]+>", " ", raw)
        return path.stem, re.sub(r"\s+", " ", stripped).strip()


def load_pdf(path: Path) -> tuple[str, str]:
    """PDF via PyMuPDF. Keeps page breaks as markdown headers so chunker can split."""
    try:
        import fitz  # type: ignore
    except ImportError as e:
        raise RuntimeError("PyMuPDF not installed: pip install pymupdf") from e

    doc = fitz.open(str(path))
    parts: list[str] = []
    title = path.stem
    for i, page in enumerate(doc):
        text = page.get_text("text")
        if not text.strip():
            continue
        parts.append(f"\n\n## Page {i + 1}\n\n{text.strip()}")
    if doc.metadata.get("title"):
        title = doc.metadata["title"].strip() or title
    doc.close()
    return title, "".join(parts).strip()


def load_docx(path: Path) -> tuple[str, str]:
    """Word .docx → markdown. Preserves heading levels + table structure.

    - Paragraphs styled Heading 1..6 become ``#`` .. ``######``.
    - Tables are rendered as pipe-style markdown tables so the chunker can
      pack them as atomic blocks.
    - Block order follows the document body, not paragraph-then-table. This
      matters for reports where tables sit inside a section.
    """
    try:
        import docx  # type: ignore
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as e:
        raise RuntimeError("python-docx not installed: pip install python-docx") from e

    doc = docx.Document(str(path))
    title = path.stem
    # Try core properties for a better title.
    if doc.core_properties.title and doc.core_properties.title.strip():
        title = doc.core_properties.title.strip()

    def _iter_block_items(parent):
        # Yield paragraphs and tables in their native document order.
        body = parent.element.body
        for child in body.iterchildren():
            if child.tag.endswith("}p"):
                yield Paragraph(child, parent)
            elif child.tag.endswith("}tbl"):
                yield Table(child, parent)

    def _render_table(tbl) -> str:
        rows: list[list[str]] = []
        for row in tbl.rows:
            rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
        if not rows:
            return ""
        header, *body = rows if len(rows) > 1 else (rows[0], [])
        out = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * len(header)) + " |"]
        for r in body:
            out.append("| " + " | ".join(r) + " |")
        return "\n".join(out)

    parts: list[str] = []
    for block in _iter_block_items(doc):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            style = (block.style.name or "").lower()
            m = re.match(r"heading\s*(\d+)", style)
            if m:
                level = min(max(int(m.group(1)), 1), 6)
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)
        else:  # Table
            md = _render_table(block)
            if md:
                parts.append(md)
    return title, "\n\n".join(parts)


def load_pptx(path: Path) -> tuple[str, str]:
    """PowerPoint .pptx → markdown. One ``## Slide N: <title>`` per slide.

    Extracts text from title placeholders, body shapes, tables, and notes.
    Speaker notes are preserved under a ``> Notes:`` block so their content
    is retrievable without being confused with slide body text.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise RuntimeError("python-pptx not installed: pip install python-pptx") from e

    prs = Presentation(str(path))
    title = path.stem
    if prs.core_properties.title and prs.core_properties.title.strip():
        title = prs.core_properties.title.strip()

    def _shape_text(shape) -> str:
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
        if getattr(shape, "has_table", False):
            rows: list[list[str]] = []
            for row in shape.table.rows:
                rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
            if not rows:
                return ""
            header, *body = rows if len(rows) > 1 else (rows[0], [])
            out = [
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(["---"] * len(header)) + " |",
            ]
            for r in body:
                out.append("| " + " | ".join(r) + " |")
            return "\n".join(out)
        return ""

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        body_lines: list[str] = []
        for shape in slide.shapes:
            s = _shape_text(shape)
            if not s:
                continue
            if shape == slide.shapes.title and shape.has_text_frame:
                slide_title = s
            else:
                body_lines.append(s)
        header = f"## Slide {i}" + (f": {slide_title}" if slide_title else "")
        parts.append(header)
        parts.extend(body_lines)
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append("> Notes: " + slide.notes_slide.notes_text_frame.text.strip())

    return title, "\n\n".join(p for p in parts if p.strip())


def load_xlsx(path: Path) -> tuple[str, str]:
    """Excel .xlsx → markdown. One ``## Sheet: <name>`` per sheet.

    Each sheet is serialised as a pipe-style markdown table. Completely empty
    rows/columns are dropped. Only the first 100 rows x 20 cols per sheet are
    emitted to keep chunk sizes reasonable.
    """
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as e:
        raise RuntimeError("openpyxl not installed: pip install openpyxl") from e

    wb = load_workbook(str(path), data_only=True, read_only=True)
    title = path.stem
    if wb.properties.title and wb.properties.title.strip():
        title = wb.properties.title.strip()

    parts: list[str] = []
    for sheet in wb.worksheets:
        rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True, max_row=100):
            vals = ["" if c is None else str(c).strip() for c in row[:20]]
            if any(v for v in vals):
                rows.append(vals)
        if not rows:
            continue
        # Truncate trailing empty columns
        width = max(len(r) for r in rows)
        for r in rows:
            while len(r) < width:
                r.append("")
        header, *body = rows if len(rows) > 1 else (rows[0], [])
        parts.append(f"## Sheet: {sheet.title}")
        parts.append("| " + " | ".join(header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for r in body:
            parts.append("| " + " | ".join(r) + " |")

    wb.close()
    return title, "\n\n".join(parts)


def load_image(path: Path) -> tuple[str, str]:
    """Image → (title, flat_markdown_of_regions).

    Runs the full preprocess → VLM layout → region extract pipeline. The
    returned text is a markdown-ish serialisation so the downstream
    chunker still works (it doesn't need to know it's dealing with an
    image). For region-typed chunks use ``rag.vision.image_to_chunks``
    directly instead of going through the generic text chunker.
    """
    try:
        from ..vision import analyse_layout_cached
    except ImportError as e:
        raise RuntimeError(f"vision module missing: {e}") from e

    payload = analyse_layout_cached(path.read_bytes())
    parts: list[str] = []
    for r in payload.get("regions") or []:
        t = (r.get("type") or "text").lower()
        if t == "title":
            parts.append(f"# {r.get('text', '')}")
        elif t == "table":
            parts.append(r.get("markdown") or r.get("text") or "")
        elif t == "figure":
            parts.append(f"*[Figure]* {r.get('text', '')}")
        elif t == "code":
            parts.append(f"```{r.get('language') or ''}\n{r.get('text', '')}\n```")
        else:
            parts.append(r.get("text") or "")

    return path.stem, "\n\n".join(p for p in parts if p.strip())


LOADERS = {
    ".md": load_markdown,
    ".markdown": load_markdown,
    ".txt": load_text,
    ".html": load_html,
    ".htm": load_html,
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".png": load_image,
    ".webp": load_image,
}


def load_file(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    loader = LOADERS.get(suffix)
    if loader is None:
        raise ValueError(f"Unsupported file type: {suffix}")
    return loader(path)


def supported_extensions() -> set[str]:
    return set(LOADERS.keys())
